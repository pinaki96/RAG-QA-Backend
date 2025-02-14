import os
import PyPDF2
import docx
import pandas as pd
import csv
from pptx import Presentation
from model.document_model import Document


def extract_text_from_file(file_path):
    """Extracts text from an uploaded document file."""
    try:
        text = ""
        file_extension = os.path.splitext(file_path)[1].lower()

        if file_extension == ".pdf":
            with open(file_path, "rb") as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
        elif file_extension in [".txt", ".md"]:
            with open(file_path, "r", encoding="utf-8") as file:
                text = file.read()
        elif file_extension == ".docx":
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file_extension in [".xls", ".xlsx"]:
            df = pd.read_excel(file_path, sheet_name=None)
            for sheet_name, sheet in df.items():
                text += f"Sheet: {sheet_name}\n"
                text += sheet.to_string(index=False) + "\n"
        elif file_extension == ".csv":
            with open(file_path, "r", encoding="utf-8") as file:
                reader = csv.reader(file)
                for row in reader:
                    text += ", ".join(row) + "\n"
        elif file_extension == ".pptx":
            ppt = Presentation(file_path)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            raise ValueError("Unsupported file type")

        return text.strip()
    except Exception as e:
        print(f"Error extracting text from file: {e}")
        return None


def create_document(name, content):
    """Handles business logic for document creation."""
    document = Document(name=name, content=content)
    return document.save()


def list_documents():
    """Fetches all documents."""
    return [doc.serialize for doc in Document.get_all()]


def get_document(doc_id):
    """Fetches a document by ID."""
    document = Document.get_by_id(doc_id)
    return document.serialize if document else None


def delete_document(doc_id):
    """Deletes a document by ID."""
    return Document.delete_by_id(doc_id)


def update_document(doc_id, name, content):
    """Updates a document and regenerates its embedding."""
    return Document.update_by_id(doc_id, name, content)
